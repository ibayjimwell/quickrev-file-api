from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import pypandoc

def convert_pdf_to_txt(file_path) -> str:

    text = ""
    try:
        reader = PdfReader(file_path)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    except Exception as e:
        print(f"Error reading PDF file {file_path}: {e}")
    return text


def convert_docx_to_txt(file_path) -> str:

    text = ""
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"Error reading DOCX file {file_path}: {e}")
    return text


def convert_pptx_to_txt(file_path) -> str:

    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading PPTX file {file_path}: {e}")
    return text


def convert_txt_to_txt(file_path) -> str:

    text = ""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()
    except Exception as e:
        print(f"Error reading TXT file {file_path}: {e}")
    return text


def convert_md_to_docx(md_file_path: str, docx_output_path: str):
    """
    Converts a Markdown file to a DOCX file using pypandoc.

    Args:
        md_file_path: The path to the input Markdown file.
        docx_output_path: The path where the output DOCX file will be saved.
    """
    try:
        pypandoc.convert_file(
            md_file_path,
            'docx',
            format='md',
            outputfile=docx_output_path
        )
    except Exception as e:
        print(f"Error converting '{md_file_path}' to DOCX: {e}")
        raise e
   
   